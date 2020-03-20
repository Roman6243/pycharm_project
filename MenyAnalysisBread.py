from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
import numpy as np
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
from matplotlib import rcParams
import os
pd.set_option('display.max.columns', 35)
pd.set_option('display.width', 1000)
rcParams['font.family'] = 'sans-serif'
rcParams['font.sans-serif'] = 'Avenir Book'
os.chdir(r'C:\Users\PC1\Dropbox (BigBlue&Company)\ETC Insight\Projects\Meny Analysis')
files = os.listdir('input data')
weekly_compagine_2018 = pd.read_excel('analyse av salgstall/Ukens brød 2019 .xlsx', sheet_name='2018_cp')
weekly_compagine_2019 = pd.read_excel('analyse av salgstall/Ukens brød 2019 .xlsx', sheet_name='2019_cp')
data = pd.DataFrame() # create empty data frame

for file in files:
    temp = pd.read_excel('input data/' + file)
    temp['Week_Year'] = temp['Week_Year'].fillna(method = 'ffill')
    temp = temp.dropna(subset = ['Leverandør'])
    meny_position = file.split('.')[0].find('Meny') # find position number of Meny
    meny_name = file.split('.')[0][meny_position::]
    temp['store'] = meny_name
    data = pd.concat([data, temp])
    print("Extract data from - %s - file" %file)

data=data.rename(columns = {'Vare': 'product', 'Leverandør':'supplier', 'Navn':'category', 'Navn.1': 'under_category', 'Navn.2' : 'category_level_3', 'Navn.3' : 'sub_category',
                            'Netto omsetning':'turnover', 'Antall solgt' : 'number_products_sold', 'Antall solgt på tilbud' : 'number_products_sold_compagine', 'Svinn kr' : 'westage'})
data[['uke_name', 'week', 'year']] = data['Week_Year'].str.split(' ', expand = True) # extract week and year
data=data.drop(columns=['category_level_3', 'uke_name']) #two the same columns category_level_3 and under_category
data['week'] = data['week'].astype(np.int64)
data['year'] = data['year'].astype(np.int64)
data['EAN/PLU'] = data['EAN/PLU'].astype(np.int64)
# there are product where 'Best.nr' == NA, but it isn't product from compagin prodacts
# set(products_in_compagin['product'].unique()).difference(set(data['product'].unique()))
data['Best.nr'] = data['Best.nr'].replace(np.nan, 0).astype(np.int64)
data.loc[(data['product'] == 'ALPEBRØD GROVT 570G')|(data['product'] == 'ALPEBRØD 570G JACOBS UTVALGTE'), 'product'] = 'ALPEBRØD HALVSTEKT 600G JACOBS' # eliminate product name "ALPEBRØD GROVT 570G"
data[['turnover', 'number_products_sold', 'number_products_sold_compagine', 'westage']] = data[['turnover', 'number_products_sold', 'number_products_sold_compagine', 'westage']].replace(np.nan, 0)
# data_bread = data[data['category'] == 'Ferske bakerivarer'].copy()

data_bread = data[data['category'] == 'Ferske bakerivarer'].reset_index(drop=True) # the same result as code above
# show which products were in compagin
products_in_compagin_2018=pd.merge(data_bread[['Best.nr', 'product', 'year', 'week']], weekly_compagine_2018, on=['Best.nr', 'year', 'week'], how='right').drop_duplicates().sort_values(by='week').\
    dropna()[['product_x', 'year', 'week']].rename(columns={'product_x':'product'})
products_in_compagin_2019=pd.merge(data_bread[['Best.nr', 'product', 'year', 'week']], weekly_compagine_2019, left_on=['product', 'year', 'week'], right_on = ['correct_name', 'year', 'week'], how='right')\
    [['correct_name', 'year', 'week']].drop_duplicates().sort_values(by=['year', 'week']).reset_index(drop=True).rename(columns={'correct_name': 'product'})
products_in_compagin = pd.concat([products_in_compagin_2018, products_in_compagin_2019]).reset_index(drop=True)

for i in range(len(products_in_compagin['product'])):
    data_bread.loc[(data_bread['product']==products_in_compagin['product'][i]) & (data_bread['week'] == products_in_compagin['week'][i]) & (data_bread['year'] == products_in_compagin['year'][i]), 'in_promotion'] = True
    print("Number of iteration: %s"%i)
data_bread.loc[data_bread['in_promotion'] != True, 'in_promotion'] = False
# show all promotions data_bread[data_bread['in_promotion']==True]
data_bread = data_bread.groupby(['product', 'category', 'under_category', 'sub_category', 'store', 'week', 'year', 'in_promotion'])\
    ['turnover', 'number_products_sold', 'number_products_sold_compagine', 'westage'].sum().reset_index()

kpi_1 = data_bread.groupby(['store', 'week', 'year', 'in_promotion'])['turnover'].sum().reset_index()
kpi_1=pd.pivot_table(kpi_1, columns = ['in_promotion', 'store'] , index = ['week', 'year'], dropna=True).reset_index()
kpi_1.columns = kpi_1.columns.droplevel().droplevel()
kpi_1.columns = ['week', 'year', 'Meny_Manglerud_False', 'Meny_Nordstrand_False', 'Meny_Oslo_City_False', 'Meny_Ringnes_Park_False', 'Meny_totalt_False',
                 'Meny_Manglerud_True', 'Meny_Nordstrand_True', 'Meny_Oslo_City_True', 'Meny_Ringnes_Park_True', 'Meny_totalt_True']
kpi_1['share_Meny_Manglerud'] = round(kpi_1['Meny_Manglerud_True']/(kpi_1['Meny_Manglerud_True']+kpi_1['Meny_Manglerud_False']), 4)
kpi_1['share_Meny_Nordstrand'] = round(kpi_1['Meny_Nordstrand_True']/(kpi_1['Meny_Nordstrand_True']+kpi_1['Meny_Nordstrand_False']), 4)
kpi_1['share_Meny_Oslo_City'] = round(kpi_1['Meny_Oslo_City_True']/(kpi_1['Meny_Oslo_City_True']+kpi_1['Meny_Oslo_City_False']), 4)
kpi_1['share_Meny_Ringnes_Park'] = round(kpi_1['Meny_Ringnes_Park_True']/(kpi_1['Meny_Ringnes_Park_True']+kpi_1['Meny_Ringnes_Park_False']), 4)
kpi_1['share_Meny_totalt'] = round(kpi_1['Meny_totalt_True']/(kpi_1['Meny_totalt_True']+kpi_1['Meny_totalt_False']), 4)
kpi_1=kpi_1.sort_values(by=['year', 'week'])
kpi_1=kpi_1[['week', 'year', 'share_Meny_Manglerud', 'share_Meny_Nordstrand', 'share_Meny_Oslo_City', 'share_Meny_Ringnes_Park', 'share_Meny_totalt']]
kpi_1.to_excel('kpi_1.xlsx')

data_bread_cat = data.copy()
data_bread_cat.loc[data_bread_cat['category']=='Ferske bakerivarer', 'bread_category'] = True
data_bread_cat.loc[data_bread_cat['category']!='Ferske bakerivarer', 'bread_category'] = False
kpi_3=data_bread_cat.groupby(['store', 'week', 'year', 'bread_category'])['turnover'].sum().reset_index()
kpi_3 = pd.pivot_table(kpi_3, columns = 'bread_category', index = ['store', 'week', 'year'], dropna=True).reset_index()
kpi_3.columns = kpi_3.columns.droplevel()
kpi_3.columns = ['store', 'week', 'year', 'bread_category_no', 'bread_category_yes']
kpi_3['share']=round(kpi_3['bread_category_yes']/(kpi_3['bread_category_no']+kpi_3['bread_category_yes']),4)
kpi_3=kpi_3.sort_values(by=['store', 'year', 'week']).dropna()


# share on the under category level
for store in data_bread['store'].unique():
    for undercategory in data_bread['under_category'].unique():
        data_bread_undercategory=data_bread[(data_bread['store'] == store) & (data_bread['under_category'] == undercategory)].groupby(by=['year', 'week', 'in_promotion'])['turnover'].sum().reset_index()
        data_bread_undercategory['total']=data_bread_undercategory.groupby(['year', 'week'])['turnover'].transform('sum')
        data_bread_undercategory['share']=data_bread_undercategory['turnover']/data_bread_undercategory['total']
        data_bread_undercategory=data_bread_undercategory[data_bread_undercategory['in_promotion']==True]
        if data_bread_undercategory.empty == False:
            print('The lenth of data_frame ' + str(store) + " " + str(undercategory) + " is " + str(len(data_bread_undercategory)))
            data_bread_undercategory.to_excel('output data/under_category_'+ undercategory.replace('/', '_') +' ' + str(store)+'.xlsx')

# point number 4
data_bread.groupby(['store', 'week', 'year', 'product'])['turnover'].sum().reset_index().sort_values(by=['year', 'week', 'turnover'], ascending=[True, True, False]).\
    groupby(['store', 'week', 'year']).head(20).sort_values(by=['year', 'week', 'store'], ascending=[True, True, False]).to_excel('data_bread_point_4.xlsx')
# products were not saled
data[data['turnover']==0].groupby(['store', 'week', 'year', 'product'])['turnover'].sum().reset_index().sort_values(by=['year', 'week', 'turnover'], ascending=[True, True, False]).\
    to_excel('products_not_saled.xlsx')

# data for plotting
plot_data = kpi_1.copy()
plot_data['x_data'] = plot_data['week'].astype(str).str.rjust(2, '0') + '/' + plot_data['year'].astype(str)
fig, ax = plt.subplots()
ax.plot('x_data', 'share_Meny_Manglerud', data=plot_data, linestyle='-', marker='o')
ax.grid()
for tick in ax.get_xticklabels():
    tick.set_fontsize(10)
    tick.set_rotation(90)
ax.set_yticklabels(['{:,.0%}'.format(x) for x in ax.get_yticks()])
ax.set_xlim(left=-1, xmax=len(plot_data))
ax.set_title('Share of products in total', fontsize=16)
ax.set_ylabel(ylabel = 'share [%]', fontsize=16)
ax.set_xlabel(xlabel = 'week/year', fontsize=16)
fig.set_size_inches(20, 10)
fig.tight_layout()
fig.savefig("plots/test.png", dpi=200, pad_inches=0.5)



prs = Presentation('C:\\Users\\PC1\\Dropbox (BigBlue&Company)\\ETC Insight\\Projects\\Meny Analysis\\presentation\\BigBlue_main_template.pptx')
# slide #1
slide = prs.slides.add_slide(prs.slide_layouts[2]) # add new slide with previous defined layout, define the number of presentation layout
title = slide.shapes.title
title.text = "List of unique \"under categories\" of \" Ferske bakerivarer\""
text_content = slide.shapes.placeholders[10]
tf = text_content.text_frame
under_cat_list = data_bread['under_category'].unique().tolist()
for under_cat in under_cat_list:
    p = tf.add_paragraph()
    p.text = under_cat
    p.level = 1
# slide #2
slide = prs.slides.add_slide(prs.slide_layouts[2]) # add new slide with previous defined layout, define the number of presentation layout
title = slide.shapes.title
title.text = "List of unique \"sub categories\" of \" Ferske bakerivarer\""
text_content = slide.shapes.placeholders[10]
tf = text_content.text_frame
sub_cat_list = data_bread['sub_category'].unique().tolist()
for sub_cat in sub_cat_list:
    p = tf.add_paragraph()
    p.text = sub_cat
    p.font.size = Pt(10)
    p.level = 1
# slide #3 and others
for compagine_product in weekly_compagine['product'].unique():
    simmilar_product = data_bread['product'][data_bread['product'].str.contains(compagine_product)].unique()

    slide = prs.slides.add_slide(prs.slide_layouts[2])  # add new slide with previous defined layout, define the number of presentation layout
    title = slide.shapes.title
    title.text = "Similar products name in database to " + compagine_product
    text_content = slide.shapes.placeholders[10]
    tf = text_content.text_frame
    sub_cat_list = data_bread['sub_category'].unique().tolist()
    for similar_product in simmilar_product:
        p = tf.add_paragraph()
        p.text = similar_product
        p.font.size = Pt(10)
        p.level = 1

prs.save('C:\\Users\\PC1\\Dropbox (BigBlue&Company)\\ETC Insight\\Projects\\Meny Analysis\\new_name.pptx')
os.startfile('C:\\Users\\PC1\\Dropbox (BigBlue&Company)\\ETC Insight\\Projects\\Meny Analysis\\new_name.pptx')

# to see placeholders name and id numbers
for shape in slide.placeholders:
	print('%d %s' % (shape.placeholder_format.idx, shape.name))
